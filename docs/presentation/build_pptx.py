"""
SpiroBird — generator prezentacije (PowerPoint .pptx)

Pokretanje:  python build_pptx.py
Izlaz:       SpiroBird-prezentacija.pptx

Slike se čitaju iz  img/  poddirektorija; ako slika ne postoji, umjesto nje se
crta placeholder okvir s uputom što uslikati. Nakon dodavanja slika samo ponovno
pokrenuti skriptu.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
OUT = os.path.join(HERE, "SpiroBird-prezentacija.pptx")

# Paleta — usklađena s igrom (tamno nebo, žuta ptica)
BG     = RGBColor(0x0B, 0x1D, 0x33)
PANEL  = RGBColor(0x12, 0x2A, 0x47)
YELLOW = RGBColor(0xF4, 0xCA, 0x40)
TEXT   = RGBColor(0xE8, 0xEE, 0xF6)
MUTED  = RGBColor(0x93, 0xA8, 0xC4)
CYAN   = RGBColor(0x41, 0xB9, 0xD5)
GREEN  = RGBColor(0x38, 0xC1, 0x72)
RED    = RGBColor(0xE3, 0x55, 0x4F)

SW, SH = Inches(13.333), Inches(7.5)


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def txt(slide, left, top, w, h, lines, size=18, color=TEXT, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for line in lines:
        if isinstance(line, str):
            line = (line, 0, color, size, False)
        text, lvl, col, sz, b = (list(line) + [0, color, size, bold])[:5]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.level = lvl
        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.color.rgb = col
        r.font.bold = b
        r.font.name = "Segoe UI"
    return box


def title(slide, text, sub=None):
    txt(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.9),
        [(text, 0, YELLOW, 32, True)])
    if sub:
        txt(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.5),
            [(sub, 0, MUTED, 16, False)])
    bar = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(2.2), Pt(3))
    bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background()


def bullets(slide, items, left=0.6, top=1.5, w=12.1, h=5.6, size=18):
    rows = []
    for it in items:
        if isinstance(it, str):
            rows.append(("• " + it, 0, TEXT, size, False))
        else:
            text, lvl = it[0], it[1] if len(it) > 1 else 0
            col = it[2] if len(it) > 2 else (TEXT if lvl == 0 else MUTED)
            b = it[3] if len(it) > 3 else False
            pre = "• " if lvl == 0 else "– "
            rows.append((pre + text, lvl, col, size - 2 * lvl, b))
    txt(slide, Inches(left), Inches(top), Inches(w), Inches(h), rows)


def image_slot(slide, fname, left, top, w, h, label):
    path = os.path.join(IMG, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 width=Inches(w), height=Inches(h))
    else:
        ph = slide.shapes.add_shape(1, Inches(left), Inches(top),
                                    Inches(w), Inches(h))
        ph.fill.solid(); ph.fill.fore_color.rgb = PANEL
        ph.line.color.rgb = CYAN; ph.line.width = Pt(1.5)
        tf = ph.text_frame; tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = f"[ SLIKA: {fname} ]\n{label}"
        r.font.size = Pt(12); r.font.color.rgb = MUTED


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    # ---- 1. Naslovna ----
    s = new_slide(prs)
    txt(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.4),
        [("SpiroBird", 0, YELLOW, 60, True)])
    txt(s, Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.0),
        [("IoT digitalni spirometar u obliku igre u stilu Flappy Bird",
          0, TEXT, 26, False)])
    txt(s, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.2), [
        ("Razvoj ugradbenih sustava · ak. god. 2025./2026. · Tehničko veleučilište u Zagrebu", 0, MUTED, 16, False),
        ("", 0, MUTED, 8, False),
        ("Hrvoje Renato Šokčić — razvoj softverskog rješenja", 0, TEXT, 16, False),
        ("Ivan Benčić — simulacija          Domagoj Lepen — testiranje", 0, TEXT, 16, False),
        ("Sven Gavranović — dokumentacija", 0, TEXT, 16, False),
    ])

    # ---- 2. Problem i motivacija ----
    s = new_slide(prs); title(s, "Problem i motivacija")
    bullets(s, [
        "Poticajna spirometrija = standardna respiratorna vježba (rehabilitacija, KOPB, postoperativno)",
        ("vježbe su monotone → pacijenti ih preskaču", 1),
        ("mehanički spirometri (kuglice) ne daju mjerljive podatke", 1),
        "Ideja: pretvoriti vježbu u igru s mjerljivim, trajno pohranjenim rezultatima",
        "Korisnik 'dahom' upravlja pticom — cilj: stabilan protok u medicinskoj zoni",
        ("pragovi 600 / 900 / 1200 ml/s preuzeti iz medicinskih smjernica", 1),
        ("uspjeh = protok 900–1200 ml/s stabilno NEPREKIDNO 5 sekundi", 1, GREEN),
        ("protok > 1200 ml/s = trenutni neuspjeh", 1, RED),
        "U prototipu dah emulira potenciometar — signalni put identičan stvarnom senzoru protoka",
    ])

    # ---- 3. Arhitektura ----
    s = new_slide(prs); title(s, "Arhitektura sustava", "dva fizička ESP32-S3 + backend")
    bullets(s, [
        ("CONTROLLER / MASTER (ESP32-S3-N16R8)", 0, YELLOW, True),
        ("jedini izvor istine: senzor, logika vježbe, aktuatori, Wi-Fi, NVS", 1),
        ("DISPLAY / SLAVE (ESP32-S3 + 2.8\" ILI9341V + kapacitivni touch)", 0, YELLOW, True),
        ("isključivo renderira primljeno — ne spaja se na Wi-Fi, nema touch izbornike", 1),
        ("ESP-NOW broadcast 40 Hz (48 B paket: magic + verzija + seq + checksum)", 0, CYAN),
        ("BACKEND: Node.js/Express + dashboard — javno na spirobird.onrender.com", 0, CYAN),
        ("HTTP POST rezultata SAMO nakon završetka pokušaja, timeout 1500 ms", 1),
        ("Offline-first: bez Wi-Fi-ja i servera igra radi potpuno normalno", 0, GREEN, True),
    ], top=1.4, h=3.4)
    image_slot(s, "arhitektura.png", 8.6, 4.6, 4.2, 2.6,
               "blok dijagram (može screenshot iz wikija, str. 4)")

    # ---- 4. Hardver ----
    s = new_slide(prs); title(s, "Hardver i ožičenje")
    bullets(s, [
        "Potenciometar B10K → ADC GPIO4 (emulator protoka daha)",
        "Vibro motor preko NPN 2N2222A: 1 kΩ baza, 10 kΩ pull-down, flyback dioda",
        "Wake gumb GPIO21 (start / skip portala / buđenje iz deep sleepa — RTC EXT0)",
        "Onboard WS2812 RGB LED — statusna indikacija stanja",
        "Napajanje 5 V / 1 A, zajednički GND; motor (3 V) na 5 V samo u kratkim pulsevima",
    ], w=6.6, size=16)
    image_slot(s, "shema-spajanja.png", 7.4, 1.4, 5.4, 3.0,
               "shema iz Circuit Designera")
    image_slot(s, "foto-setup.jpg", 7.4, 4.6, 5.4, 2.6,
               "fotografija stvarnog postava (oba uređaja)")

    # ---- 5. Obrada signala ----
    s = new_slide(prs); title(s, "Obrada signala", "ADC → protok u ml/s, 100 Hz")
    bullets(s, [
        "Uzorkovanje 100 Hz, usrednjavanje 4 uzorka, korisni ADC raspon 200–3900",
        "Kalibracija s fiksnim centrom: knob u zoni 1700–2300 držan 2 s → nul-točka = 2000",
        ("display uživo navodi korisnika: 'turn UP / turn DOWN / hold it right there!'", 1),
        "Deadzone ±70 → otklon se linearno mapira u 0–1400 ml/s",
        "EMA filter (α = 0.20) — uklanja šum bez osjetnog kašnjenja",
        "Stabilnost: ring buffer 500 ms, peak-to-peak varijacija < 180 ml/s",
        "Volumen numeričkom integracijom: V += protok × dt",
        ("izmjereni šum na hardveru: p2p < 2 ml/s — izvrsno", 1, GREEN),
    ])

    # ---- 6. State machine ----
    s = new_slide(prs); title(s, "Logika vježbe — state machine")
    bullets(s, [
        ("IDLE → CALIBRATING → READY → ACTIVE → SUCCESS / FAIL → RESULT → IDLE  (+ SLEEP)",
         0, CYAN, True),
        "Svaka tranzicija se ispisuje na serijsku konzolu (dijagnostika)",
        "ACTIVE: integracija volumena, praćenje max/prosjeka, stable timer 0→5000 ms",
        ("FAIL_OVER_1200: filtrirani > 1200 ili sirovi > 1250 (zaštita od šiljaka)", 1, RED),
        ("FAIL_TIMEOUT: pokušaj dulji od 60 s", 1),
        "Nova runda: pomak knoba (uz re-arm na mirovanju) ili gumb (puna rekalibracija)",
        "Bez blokirajućeg delay() u cijeloj logici — millis() raspoređivanje",
        "ISR-ovi postavljaju samo volatile zastavice (siguran ISR dizajn)",
    ])

    # ---- 7. ESP-NOW ----
    s = new_slide(prs); title(s, "ESP-NOW komunikacija", "bez uparivanja MAC adresa")
    bullets(s, [
        "Broadcast FF:FF:FF:FF:FF:FF, 40 Hz — Display ne treba nikakvu konfiguraciju",
        "Problem: ESP-NOW mora biti na istom kanalu kao Wi-Fi Controllera",
        ("Controller: kanal prati Wi-Fi; offline = fiksni kanal 1", 1),
        ("Display: skenira kanale 1–13 (300 ms/kanal), lock nakon 3 valjana paketa", 1),
        ("gubitak > 1 s → NO SIGNAL; > 3 s → automatski re-scan i re-lock", 1),
        "Paket 48 B: magic 0x5342, verzija, redni broj, XOR checksum",
        ("neispravni paketi se odbacuju i broje; layout čuva static_assert na obje strane", 1),
        ("Display NIKAD ne ide u deep sleep — radio bi prestao slušati", 0, YELLOW),
    ], w=7.4)
    image_slot(s, "foto-igra.jpg", 8.2, 1.4, 4.6, 5.4,
               "display s igrom uživo (ACTIVE)")

    # ---- 8. Wi-Fi provisioning ----
    s = new_slide(prs); title(s, "Wi-Fi provisioning", "captive portal — bez hardkodiranih lozinki")
    bullets(s, [
        "Spremljena mreža → spajanje 15 s → online",
        "Bez spremljene mreže → AP 'SpiroBird-Setup' + captive portal (192.168.4.1)",
        ("Display istovremeno prikazuje upute (status stiže ESP-NOW-om!)", 1),
        "Spremljena mreža nedostupna → korisnik bira na gumbu:",
        ("kratki pritisak = igraj OFFLINE   ·   dugi (1 s) = otvori portal", 1, CYAN),
        ("display crta progress bar držanja gumba uživo", 1),
        "Gumb držan pri uključivanju → portal odmah (promjena mreže)",
        "Gubitak veze u radu: max 2 pokušaja pa trajni offline (štiti ESP-NOW kanal)",
        ("OFFLINE je deterministično, terminalno stanje — demo se ne može 'zaglaviti'", 0, GREEN),
    ], w=7.4, size=16)
    image_slot(s, "foto-decision.jpg", 8.2, 1.4, 4.6, 2.6,
               "decision ekran (kratki/dugi pritisak)")
    image_slot(s, "screenshot-portal.png", 8.2, 4.2, 4.6, 2.6,
               "captive portal na mobitelu")

    # ---- 9. Igra / render ----
    s = new_slide(prs); title(s, "Igra na displayu", "LovyanGFX · 30 FPS · double buffering")
    bullets(s, [
        "Full-screen sprite 320×240 @ 8bpp (76.8 KB) → bez treperenja, 1 SPI push po frameu",
        "Ptica prati filtrirani protok s inercijom (spring + damping fizika)",
        "Linije 600/900/1200, ciljni koridor, kozmetičke cijevi, stability bar 0–5 s",
        "HUD: protok, volumen, stable timer + Wi-Fi/server status (javlja Controller)",
        "Ekrani za sva stanja: kalibracija s coachingom, uspjeh/neuspjeh, rezultat, sleep…",
        "Fake-demo mod: bez Controllera nakon 8 s vrti cijelu demo partiju (testiranje)",
        "Pseudo sleep: 1 min → dim, +15 s → backlight off; buđenje DODIROM ili aktivnošću",
    ])

    # ---- 10. Haptika ----
    s = new_slide(prs); title(s, "Haptika i statusna indikacija")
    bullets(s, [
        ("Vibracijski motor (pattern engine, non-blocking):", 0, YELLOW),
        ("upozorenje (izlazak iz zone / blizu 1200): 300 ms", 1),
        ("uspjeh: 500 ms 'win' puls   ·   neuspjeh: uzorak 3 × 600 ms", 1),
        ("sigurnost: tvrdi limit segmenta 1000 ms, cooldown 400 ms, pin LOW prva instrukcija boota, sinkroni boot self-test", 1),
        ("RGB LED po stanju:", 0, YELLOW),
        ("plava=IDLE · ljubičasta=kalibracija · žuta=READY · zelena=ACTIVE/uspjeh · crvena=FAIL", 1),
        "Buzzer eliminiran nakon testiranja (neispravna komponenta) — kod ostaje iza zastavice",
    ])

    # ---- 11. Power management ----
    s = new_slide(prs); title(s, "Upravljanje potrošnjom")
    bullets(s, [
        ("Controller:", 0, YELLOW),
        ("60 s bez inputa → pseudo sleep (iz BILO KOJEG stanja)", 1),
        ("3 min → pravi deep sleep (esp_deep_sleep_start)", 1),
        ("buđenje: gumb (RTC EXT0 — radi i iz deep sleepa), pomak knoba (debounce 50 ms)", 1),
        ("Display:", 0, YELLOW),
        ("1 min → prigušenje, +15 s → backlight ugašen", 1),
        ("buđenje: dodir ekrana (FT6336G) ili automatski na aktivnost igre", 1),
        ("deep sleep NAMJERNO izostavljen — ESP-NOW prijemnik mora uvijek slušati", 1),
        "Neaktivnost = isključivo input (pomak/gumb/događaji) — nijedno stanje ne drži uređaj budnim",
    ])
    image_slot(s, "foto-sleep.jpg", 9.2, 5.0, 3.6, 2.0,
               "sleep ekran (opcionalno)")

    # ---- 12. Backend ----
    s = new_slide(prs); title(s, "Backend i dashboard", "Node.js / Express · Render.com")
    bullets(s, [
        "REST API: /health · POST/GET/DELETE /api/results · /api/highscore",
        "Pohrana: results.json (atomični upisi); dokumentirano ograničenje free tiera",
        "Dashboard: kartice (najbolji volumen, uspješnost), tablica, auto-refresh 3 s",
        ("JAVNO U PRODUKCIJI: https://spirobird.onrender.com", 0, GREEN, True),
        "Controller POST-a samo nakon SUCCESS/FAIL — nikad tijekom igre; timeout 1500 ms + retry",
    ], w=7.0, size=16)
    image_slot(s, "screenshot-dashboard.png", 7.8, 1.4, 5.0, 5.4,
               "dashboard s pravim rezultatima")

    # ---- 13. Verifikacija ----
    s = new_slide(prs); title(s, "Verifikacija i testiranje")
    bullets(s, [
        ("Simulacija PRIJE hardvera (Python, 1:1 port logike):", 0, YELLOW),
        ("23 unit testa: EMA, mapiranje, stabilnost, volumen, uspjeh/fail, protokol (48 B + checksum)", 1),
        ("5 scenarija s ispisom tranzicija — identičan format kao firmware Serial", 1),
        ("Bring-up na hardveru strogo po koracima:", 0, YELLOW),
        ("senzor → aktuatori → backend → display → ESP-NOW → puna igra → sleep", 1),
        ("8 pronađenih i riješenih problema — svi dokumentirani (simptom→uzrok→rješenje)", 1),
        ("Statistika igranja: 29 pokušaja, 66 % uspješnosti, najbolji volumen 17.5 L", 0, GREEN),
    ])

    # ---- 14. Najzanimljiviji bugovi ----
    s = new_slide(prs); title(s, "Tri najpoučnija problema s hardvera",
                              "materijal za diskusiju")
    bullets(s, [
        ("1. Motor se vrtio neprekidno od boota", 0, RED, True),
        ("async puls + watchdog u loop() koji se ne vrti dok setup() blokira na Wi-Fi portalu", 1),
        ("pouka: sigurnosna logika ne smije ovisiti o glavnoj petlji → sinkroni self-test", 1, GREEN),
        ("2. Display u petlji hvata/gubi ESP-NOW signal", 0, RED, True),
        ("Wi-Fi auto-reconnect skenira kanale i 'vuče' ESP-NOW TX kanal sa sobom", 1),
        ("pouka: deterministička stanja umjesto pozadinskih retrya → offline = terminalno", 1, GREEN),
        ("3. TFT_eSPI StoreProhibited crash na ESP32-S3", 0, RED, True),
        ("poznati neriješeni bug biblioteke (issue #3743) — plan rizika je predvidio fallback", 1),
        ("pouka: za svaki rizik unaprijed imati izlaz → migracija na LovyanGFX u 1 commitu", 1, GREEN),
    ], size=16)

    # ---- 15. Git proces ----
    s = new_slide(prs); title(s, "Git proces i kolaboracija")
    bullets(s, [
        "35+ atomičkih commitova s opisnim porukama — povijest priča razvoj uključujući bugfixeve",
        "Stabilan main + feature grane; branch protection (PR obavezan za merge)",
        "Pull requestovi s opisima, code review i merge po ulogama članova tima",
        "Wiki: opis zadatka, analiza zahtjeva (45+ zahtjeva), use-case + sekvencijski dijagrami, arhitektura (dijagrami razreda/stanja/aktivnosti u Mermaidu)",
        "docs/: plan, ožičenje, protokol, plan testiranja, izvještaji, troubleshooting",
        "Prezentacija (ova!) committana u repo — uključujući generator (build_pptx.py)",
    ], w=7.2, size=16)
    image_slot(s, "screenshot-github.png", 8.0, 1.4, 4.8, 5.4,
               "GitHub: mergani PR-ovi / network graph")

    # ---- 16. Demo plan ----
    s = new_slide(prs); title(s, "Plan demonstracije uživo")
    bullets(s, [
        "1. Uključivanje → kalibracija s coachingom na ekranu (CENTER the knob)",
        "2. Uspješna vježba: zona 900–1200, stability bar → SUCCESS + vibracija + LED",
        "3. Rezultat na javnom dashboardu (spirobird.onrender.com) u stvarnom vremenu",
        "4. Neuspjeh: protok preko 1200 → crveno nebo, FAIL + 3 vibracije",
        "5. Otpornost: gašenje Displaya i ponovno uključenje → automatski re-lock kanala",
        "6. Sleep: 60 s mirovanja → pseudo sleep; dodir budi ekran; gumb budi Controller",
        ("rezerva: offline mod — sve radi i bez interneta (ako Wi-Fi na fakultetu zakaže)", 1, YELLOW),
    ])

    # ---- 17. Zaključak ----
    s = new_slide(prs); title(s, "Zaključak i moguće nadogradnje")
    bullets(s, [
        ("Sustav je u potpunosti implementiran i validiran na stvarnom hardveru", 0, GREEN, True),
        "Demonstrirani elementi ugradbenih sustava: ADC, prekidi + ISR zastavice, event-driven dizajn bez delay(), state machine, PWM/LEDC, NVS, watchdog guardovi, bežični protokoli, sleep modovi",
        ("Nadogradnje:", 0, YELLOW),
        ("stvarni senzor protoka daha (diferencijalni tlak) — zamjena u jednom modulu", 1),
        ("korisnički profili, trajna baza, BLE aplikacija, OTA, 3D kućište", 1),
        ("", 0),
        ("Hvala na pažnji — pitanja?", 0, YELLOW, True),
    ])

    prs.save(OUT)
    print(f"OK: {OUT}")
    missing = [f for f in [
        "arhitektura.png", "shema-spajanja.png", "foto-setup.jpg",
        "foto-igra.jpg", "foto-decision.jpg", "screenshot-portal.png",
        "foto-sleep.jpg", "screenshot-dashboard.png", "screenshot-github.png",
    ] if not os.path.exists(os.path.join(IMG, f))]
    if missing:
        print("Nedostaju slike (placeholderi umetnuti):")
        for f in missing:
            print("  img/" + f)


if __name__ == "__main__":
    build()
